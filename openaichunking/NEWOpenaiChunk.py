import os
import re
import time
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.config import Settings
import pandas as pd
from dotenv import load_dotenv
import openai
from tenacity import retry, stop_after_attempt, wait_fixed
import logging

# Configure logging
logging.basicConfig(
    filename="processing.log",
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)

def log_message(message):
    print(message)
    logging.info(message)

# Load environment variables from .env file
load_dotenv()

# Get OpenAI API key from .env
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

if not OPENAI_API_KEY:
    log_message("Error: OPENAI_API_KEY not found in .env file.")
    raise ValueError("OPENAI_API_KEY not found in .env file.")

openai.api_key = OPENAI_API_KEY

# **Utility Functions**
def sanitize_collection_name(name):
    """Sanitize collection names to meet ChromaDB requirements."""
    sanitized = re.sub(r'[^a-zA-Z0-9_-]', '_', name)
    return sanitized[:63].strip('_')

def extract_text(file_path):
    """
    Extract text from various file types: .txt, .pdf, .docx, .csv, .xlsx, .pptx.
    """
    try:
        text = ""
        if file_path.endswith(".txt"):
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        elif file_path.endswith(".pdf"):
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() or ""
        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        elif file_path.endswith(".csv"):
            df = pd.read_csv(file_path)
            text = df.to_string(index=False)
        elif file_path.endswith(".xlsx"):
            df = pd.read_excel(file_path, engine="openpyxl")
            text = df.to_string(index=False)
        elif file_path.endswith(".pptx"):
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text += paragraph.text + "\n"
        return text
    except Exception as e:
        log_message(f"Error extracting text from {file_path}: {e}")
        return ""

def chunk_text(text, chunk_size=100):
    """
    Split text into chunks of a given size.
    """
    chunks = []
    paragraphs = re.split(r'\n\s*\n', text)
    for paragraph in paragraphs:
        while len(paragraph) > chunk_size:
            chunks.append(paragraph[:chunk_size])
            paragraph = paragraph[chunk_size:]
        if paragraph.strip():  # Avoid empty chunks
            chunks.append(paragraph)
    log_message(f"Chunked text into {len(chunks)} chunks.")
    return chunks

@retry(stop=stop_after_attempt(3), wait=wait_fixed(2))
def generate_embedding(chunk):
    """
    Generate embedding for a text chunk using OpenAI's API with retry logic.
    """
    response = openai.Embedding.create(
        input=chunk,
        model="text-embedding-ada-002"
    )
    return response['data'][0]['embedding']

def embed_chunks(chunks):
    """
    Generate embeddings for text chunks using OpenAI's updated API.
    """
    embeddings = []
    for chunk in chunks:
        try:
            embedding = generate_embedding(chunk)
            embeddings.append(embedding)
        except Exception as e:
            log_message(f"Error generating embedding for chunk: {chunk[:50]} - {e}")
    log_message(f"Generated embeddings for {len(embeddings)} chunks.")
    return embeddings

def store_in_chroma(file_name, chunks, embeddings, db_path="2newopen_db"):
    """
    Store chunks and embeddings in ChromaDB.
    """
    # Initialize ChromaDB Persistent Client
    client = chromadb.PersistentClient(path=db_path)
    
    # Sanitize the collection name
    collection_name = sanitize_collection_name(os.path.splitext(os.path.basename(file_name))[0])
    collection = client.get_or_create_collection(collection_name)

    # Add chunks and embeddings to ChromaDB
    for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
        try:
            collection.add(
                ids=[f"{collection_name}_chunk_{i}"],  # Unique ID for each chunk
                metadatas=[{"source": file_name}],      # Metadata about the source
                documents=[chunk],                      # The text chunk
                embeddings=[embedding]                  # Corresponding embedding vector
            )
            log_message(f"Added chunk {i} to collection {collection_name}: {chunk[:100]}...")
        except Exception as e:
            log_message(f"Error storing chunk {i}: {e}")

    log_message(f"Data from {file_name} successfully stored in ChromaDB as {collection_name}.")

# **Folder Monitoring**
class FolderHandler(FileSystemEventHandler):
    """
    Handles file events in the monitored folder.
    """
    def __init__(self, folder_to_monitor, db_path):
        self.folder_to_monitor = folder_to_monitor
        self.db_path = db_path
        self.processed_files = {}
        self.last_file_processed_time = time.time()

    def on_modified(self, event):
        if event.is_directory:
            return
        file_path = event.src_path
        if file_path.endswith((".txt", ".pdf", ".docx", ".csv", ".xlsx", ".pptx")):
            self.process_file(file_path)

    def on_created(self, event):
        self.on_modified(event)

def process_file(self, file_path):
    try:
        last_modified_time = os.path.getmtime(file_path)
        if file_path not in self.processed_files or self.processed_files[file_path] < last_modified_time:
            self.processed_files[file_path] = last_modified_time
            self.last_file_processed_time = time.time()
            log_message(f"Processing: {file_path}")

            # Attempt to extract text
            try:
                text = extract_text(file_path)
                log_message(f"Extracted text from {os.path.basename(file_path)}.")
            except Exception as e:
                log_message(f"Error extracting text from {file_path}: {e}")
                return  # Skip further processing for this file

            # Attempt to chunk text
            try:
                chunks = chunk_text(text)
                log_message(f"Chunked text from {os.path.basename(file_path)} into {len(chunks)} chunks.")
            except Exception as e:
                log_message(f"Error chunking text from {file_path}: {e}")
                return  # Skip further processing for this file

            # Attempt to generate embeddings
            try:
                embeddings = embed_chunks(chunks)
                log_message(f"Generated embeddings for {len(embeddings)} chunks from {os.path.basename(file_path)}.")
            except Exception as e:
                log_message(f"Error generating embeddings for {file_path}: {e}")
                return  # Skip further processing for this file

            # Attempt to store in ChromaDB
            try:
                store_in_chroma(file_path, chunks, embeddings, self.db_path)
            except Exception as e:
                log_message(f"Error storing data in ChromaDB for {file_path}: {e}")
    except Exception as e:
        log_message(f"Error processing file {file_path}: {e}")


# **Main Script**
if __name__ == "__main__":
    folder_to_monitor = "C:/Users/srira/Desktop/GenAi2/openai"  # Replace with your folder path
    db_path = "2newopen_db"  # Path to store the Chroma DB

    # Ensure the folder exists
    if not os.path.exists(folder_to_monitor):
        os.makedirs(folder_to_monitor)
        log_message(f"Created folder: {folder_to_monitor}")

    # Initialize folder monitoring
    event_handler = FolderHandler(folder_to_monitor, db_path)
    observer = Observer()
    observer.schedule(event_handler, folder_to_monitor, recursive=True)

    log_message(f"Monitoring folder: {folder_to_monitor}")

    try:
        observer.start()
        while True:
            time.sleep(15)  # Print log every 15 seconds
            time_since_last_file = time.time() - event_handler.last_file_processed_time
            if time_since_last_file >= 15:
                log_message("No new files have been added recently. Waiting for new files to process...")
    except KeyboardInterrupt:
        observer.stop()
    observer.join()
