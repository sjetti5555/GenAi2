"""
QARagver1: Basic Q&A System with Folder Monitoring

Description:
This Python script implements a foundational Q&A system with the following features:

1. **Folder Monitoring:**
   - Monitors a specified folder for new files in real-time.
   - Processes various file types (e.g., .txt, .pdf, .docx, .csv, .xlsx, .pptx).
   - Extracts content and stores embeddings in ChromaDB.

2. **Question-Answering Agent:**
   - Retrieves relevant chunks from ChromaDB.
   - Generates answers using OpenAI's GPT model.

3. **Interactive QA Loop:**
   - Users can ask questions and receive answers via the command line.

4. **Basic File Processing:**
   - Extracts text from supported file formats.
   - Splits text into chunks and stores them in ChromaDB.

Instructions:
- Place the script in a folder and specify the folder to monitor in `folder_to_monitor`.
- Run the script, and interact with the Q&A system via the command line.
"""

import os
import time
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain.text_splitter import CharacterTextSplitter
from langchain.chains import RetrievalQA
from langchain_community.chat_models import ChatOpenAI
from langchain_huggingface import HuggingFaceEmbeddings
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import chardet
import pdfplumber

# **Initialize Embedding Model**
embedding_model = SentenceTransformerEmbeddings(model_name="all-MiniLM-L6-v2")

# **Initialize ChromaDB**
vectorstore = Chroma(persist_directory="scalable1_agent_db", embedding_function=embedding_model)

# **Utility Functions**
def extract_text(file_path):
    """
    Extract text from various file types: .txt, .pdf, .docx, .csv, .xlsx, .pptx.
    """
    try:
        text = ""
        if file_path.endswith(".txt"):
            # Detect encoding
            with open(file_path, "rb") as f:
                raw_data = f.read()
                result = chardet.detect(raw_data)
                encoding = result["encoding"] or "utf-8"
            with open(file_path, "r", encoding=encoding, errors="ignore") as f:
                text = f.read()
        elif file_path.endswith(".pdf"):
            try:
                # Try PyPDF2 first
                reader = PdfReader(file_path)
                for page in reader.pages:
                    text += page.extract_text() or ""
            except Exception:
                # Fallback to pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() or ""
        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        elif file_path.endswith(".csv"):
            df = pd.read_csv(file_path)
            text = df.to_string(index=False)
        elif file_path.endswith(".xlsx"):
            df = pd.read_excel(file_path, engine="openpyxl")
            text = df.to_string(index=False)
        elif file_path.endswith(".pptx"):
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text += paragraph.text + "\n"
        else:
            raise ValueError(f"Unsupported file type: {file_path}")
        return text
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return ""

def process_file(file_path, vectorstore):
    """
    Process the file: extract text, chunk it, generate embeddings, and store in ChromaDB.
    """
    try:
        # Extract text from the file
        text = extract_text(file_path)
        if not text.strip():
            print(f"No text extracted from {file_path}. Skipping.")
            return

        # Split text into chunks
        splitter = CharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        chunks = splitter.split_text(text)

        # Add chunks to ChromaDB
        vectorstore.add_texts(chunks, metadatas=[{"source": file_path}] * len(chunks))
        vectorstore.persist()
        print(f"File processed and stored: {file_path}")
    except Exception as e:
        print(f"Error processing file {file_path}: {e}")

# **Folder Monitoring**
class FolderHandler(FileSystemEventHandler):
    """
    Monitor a folder and process new files added to it.
    """
    def __init__(self, vectorstore):
        self.vectorstore = vectorstore

    def on_created(self, event):
        if event.is_directory:
            return
        file_path = event.src_path
        if file_path.endswith((".txt", ".pdf", ".docx", ".csv", ".xlsx", ".pptx")):
            process_file(file_path, self.vectorstore)

# **Question Answering Agent**
def create_qa_agent(vectorstore):
    """
    Create a question-answering agent using LangChain.
    """
    retriever = vectorstore.as_retriever()
    llm = ChatOpenAI(model="gpt-3.5-turbo", temperature=0)  # Replace with "gpt-3.5-turbo" if needed
    qa_chain = RetrievalQA.from_chain_type(llm=llm, retriever=retriever)
    return qa_chain

# **Main Script**
if __name__ == "__main__":
    folder_to_monitor = "C:/Users/srira/Desktop/GenAi2/QA_Agents/qa_files"

    # Ensure the folder exists
    if not os.path.exists(folder_to_monitor):
        os.makedirs(folder_to_monitor)

    # Initialize folder monitoring
    handler = FolderHandler(vectorstore)
    observer = Observer()
    observer.schedule(handler, folder_to_monitor, recursive=True)

    print(f"Monitoring folder: {folder_to_monitor}")

    # Create QA agent
    qa_agent = create_qa_agent(vectorstore)

    try:
        # Start folder monitoring
        observer.start()
        while True:
            # Allow user to ask questions
            question = input("\nAsk a question (or type 'exit' to quit): ")
            if question.lower() == "exit":
                break
            answer = qa_agent.run(question)
            print(f"Answer: {answer}")
    except KeyboardInterrupt:
        observer.stop()
    observer.join()
