import os
import logging
import hashlib
import time
from threading import Thread, Event
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain.chains import RetrievalQA
from langchain_community.chat_models import ChatOpenAI
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Logging Configuration
logging.basicConfig(
    filename="app_log.log",
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)

def log_message(message):
    """Logs messages to console and file."""
    print(f"{message}")
    logging.info(message)

# Embedding Model Initialization
embedding_model = SentenceTransformerEmbeddings(model_name="all-MiniLM-L6-v2")

# Vectorstore Initialization
vectorstore = Chroma(persist_directory="vectorstore_db", embedding_function=embedding_model)

# Utility Functions
def extract_data(file_path):
    """Extracts data from various file types."""
    try:
        if file_path.endswith(".csv"):
            import pandas as pd
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        elif file_path.endswith(".xlsx"):
            import pandas as pd
            df = pd.read_excel(file_path, engine="openpyxl")
            return df.to_string(index=False)
        elif file_path.endswith(".txt"):
            with open(file_path, "r", encoding="utf-8") as file:
                return file.read()
        elif file_path.endswith(".pdf"):
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            text = " ".join(page.extract_text() for page in reader.pages)
            return text
        elif file_path.endswith(".docx"):
            from docx import Document
            doc = Document(file_path)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)
        elif file_path.endswith(".pptx"):
            from pptx import Presentation
            presentation = Presentation(file_path)
            slides_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slides_text.append(shape.text)
            return "\n".join(slides_text)
        else:
            raise ValueError(f"Unsupported file type: {file_path}")
    except Exception as e:
        log_message(f"Error extracting data from {file_path}: {e}")
        return ""

def process_file(file_path, vectorstore, processed_files):
    """Processes a file by extracting text, chunking it, and storing embeddings."""
    try:
        file_hash = get_file_hash(file_path)
        if processed_files.get(file_path) == file_hash:
            log_message(f"No changes detected for {file_path}. Skipping.")
            return

        text = extract_data(file_path)
        if not text.strip():
            log_message(f"No valid content in {file_path}. Skipping.")
            return

        # Chunk text
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000, chunk_overlap=100, separators=["\n\n", "\n", ".", " "]
        )
        chunks = splitter.split_text(text)

        # Add to vectorstore
        ids = [f"{file_path}_{i}" for i in range(len(chunks))]
        vectorstore.add_texts(chunks, metadatas=[{"source": file_path}] * len(chunks), ids=ids)
        vectorstore.persist()

        log_message(f"Processed and stored content from {file_path}.")
        processed_files[file_path] = file_hash
    except Exception as e:
        log_message(f"Error processing {file_path}: {e}")

def get_file_hash(file_path):
    """Generates a hash for a file."""
    with open(file_path, "rb") as f:
        return hashlib.md5(f.read()).hexdigest()

# Folder Monitoring
class FolderMonitorHandler(FileSystemEventHandler):
    """Handles folder events (create, modify)."""
    def __init__(self, vectorstore, processed_files, activity_event):
        self.vectorstore = vectorstore
        self.processed_files = processed_files
        self.activity_event = activity_event

    def on_created(self, event):
        if not event.is_directory:
            log_message(f"New file detected: {event.src_path}")
            process_file(event.src_path, self.vectorstore, self.processed_files)
            self.activity_event.set()

    def on_modified(self, event):
        if not event.is_directory:
            log_message(f"File modified: {event.src_path}")
            process_file(event.src_path, self.vectorstore, self.processed_files)
            self.activity_event.set()

def start_folder_monitoring(folder_path, vectorstore, activity_event):
    """Starts monitoring a folder in a background thread."""
    processed_files = {}
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)
        log_message(f"Created folder: {folder_path}")

    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)
        if os.path.isfile(file_path):
            process_file(file_path, vectorstore, processed_files)

    event_handler = FolderMonitorHandler(vectorstore, processed_files, activity_event)
    observer = Observer()
    observer.schedule(event_handler, folder_path, recursive=True)
    log_message(f"Monitoring folder: {folder_path}")
    observer.start()
    return observer

# Query Handling Functions
def handle_casual_question(query):
    """Handles casual queries like greetings or small talk."""
    casual_responses = {
        "hi": "Hello! How can I assist you today?",
        "hello": "Hi there! How can I help you?",
        "how are you?": "I'm here and ready to assist with any queries you have!",
        "can you help me?": "Of course! Let me know what you need help with."
    }
    return casual_responses.get(query.lower(), None)

def list_sources(vectorstore):
    """Lists all sources in the vectorstore."""
    docs = vectorstore.collection.get(include=["metadatas"])
    sources = {doc["source"] for doc in docs["metadatas"] if "source" in doc}
    if not sources:
        return "No sources found in the database."
    return "\n".join(f"- {source}" for source in sources)

def handle_database_question(query, qa_agent):
    """Handles database-related queries."""
    try:
        response = qa_agent({"query": query})

        # If no relevant sources are found
        if not response.get("source_documents"):
            return "I don't know the answer to that.", None

        # Format the response and sources
        formatted_answer = response["result"]
        sources = []
        for doc in response["source_documents"][:3]:  # Limit to top 4 sources
            snippet = doc.page_content[:200] + "..."  # Snippet of the chunk
            sources.append(f"{doc.metadata['source']}\n - \n{snippet}\n\n")

        return formatted_answer, sources
    except Exception as e:
        return f"Error retrieving answer: {e}", None

def review_output(answer, sources):
    """Reviews and formats the output for display."""
    if not answer:
        return "No relevant information found."

    formatted_output = f"Answer:\n{answer}\n\nSources:\n"
    if sources:
        for idx, source in enumerate(sources, start=1):
            formatted_output += f"{idx}. {source}\n"
    else:
        formatted_output += "None"
    return formatted_output

def query_router(query, qa_agent, vectorstore):
    """Routes queries to the appropriate handler."""
    # Handle casual questions
    casual_response = handle_casual_question(query)
    if casual_response:
        return casual_response

    # Handle source listing
    if query.lower() in ["what are the sources you have?", "list sources"]:
        sources_list = list_sources(vectorstore)
        return f"Available Sources:\n{sources_list}"

    # Handle database-related questions
    answer, sources = handle_database_question(query, qa_agent)
    return review_output(answer, sources)

# Main Script
if __name__ == "__main__":
    folder_to_monitor = "C:/Users/srira/Desktop/GenAi2/QA_Agents/qa_files"
    activity_event = Event()

    observer = start_folder_monitoring(folder_to_monitor, vectorstore, activity_event)

    try:
        retriever = vectorstore.as_retriever()
        retriever.search_kwargs = {"k": 5}
        llm = ChatOpenAI(model="gpt-3.5-turbo", temperature=0)
        qa_agent = RetrievalQA.from_chain_type(llm=llm, retriever=retriever, return_source_documents=True)

        print("Start asking questions! Type 'exit' to quit.")
        while True:
            # Wait for 1 seconds after file activity
            activity_event.wait(timeout=1)
            activity_event.clear()

            # Ask the user for input
            query = input("\nEnter your question: ").strip()
            if query.lower() == "exit":
                print("\nExiting the system. Goodbye!")
                break

            response = query_router(query, qa_agent, vectorstore)
            print("\n" + response)
    except KeyboardInterrupt:
        log_message("\nShutting down the application.")
        observer.stop()
