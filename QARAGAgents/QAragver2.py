import os
import logging
import hashlib
from concurrent.futures import ThreadPoolExecutor
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import RetrievalQA
from langchain_community.chat_models import ChatOpenAI
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import chardet
import pdfplumber
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import time
from threading import Thread

# **Logging Configuration**
logging.basicConfig(
    filename="processing.log",
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)

def log_message(message):
    print(f"[INFO] {message}")
    logging.info(message)

# **Initialize Embedding Model**
embedding_model = SentenceTransformerEmbeddings(model_name="all-MiniLM-L6-v2")

# **Initialize ChromaDB**
vectorstore = Chroma(persist_directory="scalable_agent_db", embedding_function=embedding_model)

# **Utility Functions**
def extract_text(file_path):
    """Extract text from various file types."""
    try:
        text = ""
        if file_path.endswith(".txt"):
            with open(file_path, "rb") as f:
                raw_data = f.read()
                result = chardet.detect(raw_data)
                encoding = result["encoding"] or "utf-8"
            with open(file_path, "r", encoding=encoding, errors="ignore") as f:
                text = f.read()
        elif file_path.endswith(".pdf"):
            try:
                reader = PdfReader(file_path)
                for page in reader.pages:
                    text += page.extract_text() or ""
            except Exception:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() or ""
        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        elif file_path.endswith(".csv"):
            df = pd.read_csv(file_path)
            text = df.to_string(index=False)
        elif file_path.endswith(".xlsx"):
            df = pd.read_excel(file_path, engine="openpyxl")
            text = df.to_string(index=False)
        elif file_path.endswith(".pptx"):
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text += paragraph.text + "\n"
        else:
            raise ValueError(f"Unsupported file type: {file_path}")
        return text
    except Exception as e:
        log_message(f"Error extracting text from {file_path}: {e}")
        return ""

def get_file_hash(file_path):
    """Generate a hash for the file content."""
    with open(file_path, "rb") as f:
        file_data = f.read()
    return hashlib.md5(file_data).hexdigest()

def process_file(file_path, vectorstore, processed_files):
    """Process the file: extract text, chunk it, and store in ChromaDB."""
    try:
        # Check if the file content has changed using its hash
        file_hash = get_file_hash(file_path)
        if processed_files.get(file_path) == file_hash:
            log_message(f"No changes detected for {file_path}. Skipping.")
            return

        text = extract_text(file_path)
        if not text.strip():
            log_message(f"No text extracted from {file_path}. Skipping.")
            return

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=2500,
            chunk_overlap=180,
            separators=["\n\n", "\n", " ", ""]
        )
        chunks = splitter.split_text(text)

        if not chunks:
            log_message(f"No chunks created for {file_path}. Skipping.")
            return

        # Add chunks with unique IDs to avoid duplication
        new_chunks = []
        ids = []
        for idx, chunk in enumerate(chunks):
            chunk_id = f"{file_path}_{idx}"
            ids.append(chunk_id)
            new_chunks.append(chunk)

        vectorstore.add_texts(new_chunks, metadatas=[{"source": file_path}] * len(new_chunks), ids=ids)
        log_message(f"New chunks added to vectorstore for {file_path}.")
        vectorstore.persist()
        log_message("Persisted vectorstore.")

        # Update processed files with the current hash
        processed_files[file_path] = file_hash
    except Exception as e:
        log_message(f"Error processing file {file_path}: {e}")

# **Folder Monitoring**
class FolderMonitorHandler(FileSystemEventHandler):
    """Monitor folder for file additions and modifications."""
    def __init__(self, vectorstore, processed_files):
        self.vectorstore = vectorstore
        self.processed_files = processed_files

    def on_created(self, event):
        if not event.is_directory:
            log_message(f"New file detected: {event.src_path}")
            process_file(event.src_path, self.vectorstore, self.processed_files)

    def on_modified(self, event):
        if not event.is_directory:
            log_message(f"File modified: {event.src_path}")
            process_file(event.src_path, self.vectorstore, self.processed_files)

def process_existing_files(folder_path, vectorstore, processed_files):
    """Process all existing files in the folder on startup."""
    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)
        if os.path.isfile(file_path):
            process_file(file_path, vectorstore, processed_files)

def start_folder_monitoring(folder_path, vectorstore, processed_files):
    """Start monitoring the specified folder."""
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)
        log_message(f"Created folder: {folder_path}")

    # Process existing files before starting monitoring
    log_message("Processing existing files in the folder...")
    process_existing_files(folder_path, vectorstore, processed_files)

    event_handler = FolderMonitorHandler(vectorstore, processed_files)
    observer = Observer()
    observer.schedule(event_handler, folder_path, recursive=True)

    log_message(f"Monitoring folder: {folder_path}")
    observer.start()
    try:
        observer.join()  # Keeps the monitoring thread running
    except KeyboardInterrupt:
        log_message("Stopping folder monitoring...")
        observer.stop()

# **Question Answering Agent**
def create_qa_agent(vectorstore):
    """Create a question-answering agent."""
    retriever = vectorstore.as_retriever()
    retriever.search_kwargs = {"k": 5}
    llm = ChatOpenAI(model="gpt-3.5-turbo", temperature=0)
    qa_chain = RetrievalQA.from_chain_type(llm=llm, retriever=retriever, return_source_documents=True)
    return qa_chain

def qa_loop(qa_agent):
    """Interactive QA loop."""
    print("\nYou can now ask questions (type 'exit' to quit):")
    try:
        while True:
            question = input("\nAsk a question: ")
            if question.lower() == "exit":
                print("\nExiting QA loop.")
                break
            
            # Skip retriever for generic greetings
            if question.lower() in ["hi", "hello", "hey"]:
                print("\nAnswer:")
                print("Hello! How can I assist you today?")
                continue
            
            result = qa_agent({"query": question})
            print("\nAnswer:")
            print(result["result"])
            
            # Display sources only if documents are retrieved
            if result.get("source_documents"):
                print("\nSources:")
                for doc in result["source_documents"]:
                    print(f"- {doc.metadata['source']}")
            else:
                print("\nNo relevant sources found.")
    except KeyboardInterrupt:
        print("\nExiting QA loop.")

# **Main Script**
if __name__ == "__main__":
    folder_to_monitor = "C:/Users/srira/Desktop/GenAi2/QA_Agents/qa_files"
    processed_files = {}  # Dictionary to store processed file hashes

    # Start folder monitoring in the background
    monitoring_thread = Thread(target=start_folder_monitoring, args=(folder_to_monitor, vectorstore, processed_files))
    monitoring_thread.daemon = True
    monitoring_thread.start()

    # Initialize QA agent
    qa_agent = create_qa_agent(vectorstore)
    qa_loop(qa_agent)
