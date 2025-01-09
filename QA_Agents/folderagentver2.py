import os
import time
import hashlib
from concurrent.futures import ThreadPoolExecutor
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from dotenv import load_dotenv
import chromadb
from chromadb.config import Settings
from chromadb.utils import embedding_functions
from langchain.text_splitter import RecursiveCharacterTextSplitter
from PyPDF2 import PdfReader
import pandas as pd
from openpyxl import load_workbook
import json
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import xml.etree.ElementTree as ET

# Load environment variables
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

# Initialize ChromaDB
chroma_client = chromadb.Client(Settings(
    persist_directory="./chroma_db",
))
embedding_fn = embedding_functions.OpenAIEmbeddingFunction(api_key=OPENAI_API_KEY)


class FolderMonitorAgent(FileSystemEventHandler):
    """
    Monitor a folder and process files when added or modified.
    """
    def __init__(self, folder_to_monitor):
        self.folder_to_monitor = folder_to_monitor
        self.last_processed = {}
        self.last_activity_time = time.time()

    def get_file_hash(self, file_path):
        """Generate a hash for file content and modification time."""
        try:
            file_mod_time = os.path.getmtime(file_path)
            with open(file_path, "rb") as f:
                file_content = f.read()
            hash_input = file_content + str(file_mod_time).encode()
            return hashlib.md5(hash_input).hexdigest()
        except Exception as e:
            print(f"Error hashing file {file_path}: {e}")
            return None

    def process_file(self, file_path):
        """Process a file: read, split, embed, and store."""
        if not os.path.exists(file_path):
            return

        # Detect duplicates based on hash
        file_hash = self.get_file_hash(file_path)
        if file_path in self.last_processed and self.last_processed[file_path] == file_hash:
            print(f"No changes detected for: {file_path}")
            return
        self.last_processed[file_path] = file_hash
        self.last_activity_time = time.time()

        print(f"Processing file: {file_path}")
        file_name = os.path.basename(file_path)
        collection_name = self.sanitize_collection_name(file_name.split('.')[0])

        # Ensure collection exists
        if collection_name not in [col.name for col in chroma_client.list_collections()]:
            chroma_client.create_collection(name=collection_name, embedding_function=embedding_fn)

        collection = chroma_client.get_collection(name=collection_name)

        # Read and split content
        content = self.read_file(file_path)
        if not content:
            print(f"Failed to read content from {file_path}")
            return

        chunks = self.chunk_text(content)
        self.process_chunks_multithreaded(chunks, collection, file_name)

    def sanitize_collection_name(self, name):
        """Sanitize collection name for ChromaDB."""
        import re
        sanitized = re.sub(r"[^a-zA-Z0-9_-]", "_", name)
        return sanitized[:63]

    def read_file(self, file_path):
        """Read file content based on type."""
        try:
            if file_path.endswith(".txt") or file_path.endswith(".log") or file_path.endswith(".md"):
                with open(file_path, "r", encoding="utf-8") as f:
                    return f.read()
            elif file_path.endswith(".pdf"):
                reader = PdfReader(file_path)
                return " ".join(page.extract_text() for page in reader.pages if page.extract_text())
            elif file_path.endswith(".csv"):
                df = pd.read_csv(file_path)
                return df.to_string()
            elif file_path.endswith(".json"):
                with open(file_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                return json.dumps(data, indent=4)
            elif file_path.endswith(".xlsx"):
                wb = load_workbook(file_path)
                content = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        content.append(" ".join(map(str, row)))
                return "\n".join(content)
            elif file_path.endswith(".docx"):
                doc = Document(file_path)
                return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            elif file_path.endswith(".pptx"):
                prs = Presentation(file_path)
                content = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            content.append(shape.text)
                return "\n".join(content)
            elif file_path.endswith(".xml"):
                tree = ET.parse(file_path)
                root = tree.getroot()
                return ET.tostring(root, encoding="utf-8", method="text").decode("utf-8")
            elif file_path.endswith(".html"):
                with open(file_path, "r", encoding="utf-8") as f:
                    soup = BeautifulSoup(f, "html.parser")
                return soup.get_text()
            else:
                print(f"Unsupported file type: {file_path}")
                return None
        except Exception as e:
            print(f"Error reading file {file_path}: {e}")
            return None

    def chunk_text(self, text):
        """Chunk text with context overlap."""
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=300  # Increased overlap for better context
        )
        return text_splitter.split_text(text)

    def process_chunks_multithreaded(self, chunks, collection, file_name):
        """Embed and store chunks using multithreading."""
        def process_single_chunk(idx, chunk):
            chunk_id = hashlib.md5(chunk.encode()).hexdigest()
            existing_docs = collection.get(ids=[chunk_id], include=["ids"])["ids"]
            if chunk_id not in existing_docs:
                collection.add(
                    ids=[chunk_id],
                    documents=[chunk],
                    metadatas={"file_name": file_name, "chunk_index": idx}
                )
            else:
                print(f"Duplicate chunk detected, skipping: {chunk_id}")

        with ThreadPoolExecutor(max_workers=4) as executor:
            for idx, chunk in enumerate(chunks):
                executor.submit(process_single_chunk, idx, chunk)

    def on_created(self, event):
        """Handle new file creation."""
        if not event.is_directory:
            self.process_file(event.src_path)

    def on_modified(self, event):
        """Handle file modifications."""
        if not event.is_directory:
            self.process_file(event.src_path)


def start_folder_monitoring(folder_to_monitor):
    """Monitor folder for changes."""
    if not os.path.exists(folder_to_monitor):
        os.makedirs(folder_to_monitor)
        print(f"Created folder: {folder_to_monitor}")

    event_handler = FolderMonitorAgent(folder_to_monitor)
    observer = Observer()
    observer.schedule(event_handler, folder_to_monitor, recursive=True)

    print(f"Monitoring folder: {folder_to_monitor}")
    observer.start()
    try:
        while True:
            time.sleep(10)
            time_since_last_activity = time.time() - event_handler.last_activity_time
            if time_since_last_activity > 30:
                print("No files added or modified in the last 30 seconds.")
                event_handler.last_activity_time = time.time()
    except KeyboardInterrupt:
        print("Stopping folder monitoring...")
        observer.stop()
    observer.join()


# Main Script
if __name__ == "__main__":
    folder_to_monitor = "C:/Users/srira/Desktop/GenAi2/QARAGAgents/qa_files"
    start_folder_monitoring(folder_to_monitor)
